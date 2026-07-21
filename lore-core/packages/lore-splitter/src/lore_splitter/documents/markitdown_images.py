from __future__ import annotations

import hashlib
from io import BytesIO

from docx import Document
from PIL import Image
from pptx import Presentation

from lore_splitter.documents.contracts import (
    DocumentInputArtifact,
    ImageSourceLocation,
    ImageToastCandidate,
)


class MarkItDownImageCollector:
    def __init__(self) -> None:
        self._by_file_id: dict[str, tuple[ImageToastCandidate, ...]] = {}

    def collect(self, document: DocumentInputArtifact) -> tuple[ImageToastCandidate, ...]:
        extension = document.normalized_extension.lower()
        if extension == ".docx":
            candidates = tuple(_docx_candidates(document))
        elif extension == ".pptx":
            candidates = tuple(_pptx_candidates(document))
        else:
            candidates = ()
        self._by_file_id[document.file_id] = candidates
        return candidates

    def candidates_for(self, document: DocumentInputArtifact) -> tuple[ImageToastCandidate, ...]:
        return self._by_file_id.get(document.file_id, ())


def markitdown_image_converter_factory() -> MarkItDownImageCollector:
    return MarkItDownImageCollector()


def extract_collected_markitdown_images(
    document: DocumentInputArtifact,
    collector: MarkItDownImageCollector,
) -> tuple[ImageToastCandidate, ...]:
    return collector.candidates_for(document)


def _docx_candidates(document: DocumentInputArtifact) -> list[ImageToastCandidate]:
    doc = Document(document.local_path)
    candidates: list[ImageToastCandidate] = []
    for relationship_id, rel in sorted(doc.part.rels.items()):
        target_part = getattr(rel, "target_part", None)
        content_type = getattr(target_part, "content_type", "")
        if not content_type.startswith("image/"):
            continue
        payload = target_part.blob
        width, height = _image_size(payload)
        extension = _extension_from_content_type(content_type)
        candidates.append(
            _candidate(
                document,
                payload=payload,
                content_type=content_type,
                extension=extension,
                width=width,
                height=height,
                location=ImageSourceLocation(
                    source_format="docx",
                    relationship_id=relationship_id,
                ),
            )
        )
    return candidates


def _pptx_candidates(document: DocumentInputArtifact) -> list[ImageToastCandidate]:
    presentation = Presentation(document.local_path)
    candidates: list[ImageToastCandidate] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            image = getattr(shape, "image", None)
            if image is None:
                continue
            payload = image.blob
            width, height = _image_size(payload)
            extension = (
                f".{image.ext}"
                if image.ext
                else _extension_from_content_type(image.content_type)
            )
            candidates.append(
                _candidate(
                    document,
                    payload=payload,
                    content_type=image.content_type,
                    extension=extension,
                    width=width,
                    height=height,
                    location=ImageSourceLocation(
                        source_format="pptx",
                        slide_number=slide_index,
                        shape_id=getattr(shape, "shape_id", None),
                        metadata={"shape_name": getattr(shape, "name", "")},
                    ),
                )
            )
    return candidates


def _candidate(
    document: DocumentInputArtifact,
    *,
    payload: bytes,
    content_type: str,
    extension: str,
    width: int | None,
    height: int | None,
    location: ImageSourceLocation,
) -> ImageToastCandidate:
    return ImageToastCandidate(
        payload=payload,
        content_type=content_type.lower(),
        extension=extension,
        byte_size=len(payload),
        checksum_sha256=hashlib.sha256(payload).hexdigest(),
        width_px=width,
        height_px=height,
        source_identity=document.source_identity,
        source_location=location,
    )


def _image_size(payload: bytes) -> tuple[int | None, int | None]:
    with Image.open(BytesIO(payload)) as image:
        return image.size


def _extension_from_content_type(content_type: str) -> str:
    return {
        "image/jpeg": ".jpg",
        "image/jpg": ".jpg",
        "image/png": ".png",
        "image/gif": ".gif",
        "image/webp": ".webp",
    }.get(content_type.lower(), f".{content_type.rsplit('/', 1)[-1].lower()}")
