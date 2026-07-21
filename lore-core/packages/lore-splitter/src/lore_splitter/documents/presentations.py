"""Pure PPTX extraction and bounded presentation chunk composition."""

from __future__ import annotations

import hashlib
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from lore_splitter.chunks import (
    Chunk,
    ChunkBudget,
    ChunkCoordinates,
    build_chunk,
)


@dataclass(frozen=True)
class PresentationConfig:
    max_slides_per_chunk: int = 4
    max_blocks_per_slide: int = 64
    image_max_bytes: int = 5 * 1024 * 1024


@dataclass(frozen=True)
class PresentationBlock:
    kind: str
    text: str
    x: int | None = None
    y: int | None = None
    width: int | None = None
    height: int | None = None
    payload_id: str | None = None


@dataclass(frozen=True)
class PresentationSlide:
    number: int
    title: str
    section: str
    blocks: tuple[PresentationBlock, ...]
    notes: str = ""
    deck_title: str = ""


@dataclass(frozen=True)
class PresentationResult:
    slides: tuple[PresentationSlide, ...]
    diagnostics: tuple[dict[str, Any], ...] = ()


@dataclass(frozen=True)
class PresentationChunkResult:
    chunks: tuple[Chunk, ...]
    payload_plans: tuple[Any, ...] = ()
    diagnostics: tuple[dict[str, Any], ...] = ()


def extract_presentation(
    path: str | Path, *, config: PresentationConfig | None = None
) -> PresentationResult:
    """Extract visible PPTX semantics without invoking Airflow or network services."""
    config = config or PresentationConfig()
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - dependency is project-declared.
        raise RuntimeError("presentation_dependency_missing") from exc
    presentation = Presentation(str(path))
    deck_title = getattr(presentation.core_properties, "title", "") or Path(path).stem
    slides: list[PresentationSlide] = []
    diagnostics: list[dict[str, Any]] = []
    for number, slide in enumerate(presentation.slides, 1):
        shapes = list(slide.shapes)
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape is not None else ""
        blocks: list[PresentationBlock] = []
        for index, shape in enumerate(shapes):
            if shape is title_shape or not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if not text:
                continue
            try:
                x, y, width, height = shape.left, shape.top, shape.width, shape.height
            except Exception:  # noqa: BLE001 - malformed shape coordinates are non-fatal.
                x = y = width = height = None
            kind = "table" if getattr(shape, "has_table", False) else "text"
            payload_id = None
            if kind == "table":
                payload_id = _table_id(shape)
            blocks.append(PresentationBlock(kind, text, x, y, width, height, payload_id))
        blocks.sort(
            key=lambda block: (block.y is None, block.y or 0, block.x is None, block.x or 0)
        )
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:  # noqa: BLE001 - notes are optional in python-pptx.
            pass
        section = _section_for_slide(title, number)
        slides.append(
            PresentationSlide(
                number,
                title,
                section,
                tuple(blocks[: config.max_blocks_per_slide]),
                notes,
                deck_title,
            )
        )
    return PresentationResult(tuple(slides), tuple(diagnostics))


def build_presentation_chunks(
    *,
    run_id: str,
    file_id: str,
    presentation: PresentationResult | str | Path,
    budget: ChunkBudget | None = None,
    config: PresentationConfig | None = None,
) -> PresentationChunkResult:
    budget = budget or ChunkBudget()
    config = config or PresentationConfig()
    result = (
        extract_presentation(presentation, config=config)
        if isinstance(presentation, (str, Path))
        else presentation
    )
    chunks: list[Chunk] = []
    diagnostics = list(result.diagnostics)
    ordinal = 0
    index = 0
    while index < len(result.slides):
        first = result.slides[index]
        group = [first]
        index += 1
        while index < len(result.slides) and len(group) < config.max_slides_per_chunk:
            candidate = result.slides[index]
            if candidate.section != first.section or candidate.title.lower().startswith(
                ("agenda", "contents")
            ):
                break
            if len(_slide_text(group + [candidate])) > budget.max_fulltext_chars:
                break
            group.append(candidate)
            index += 1
        built = _build_group(run_id, file_id, group, ordinal, budget)
        chunks.extend(built)
        ordinal += len(built)
    return PresentationChunkResult(tuple(chunks), (), tuple(diagnostics))


def _build_group(
    run_id: str, file_id: str, slides: list[PresentationSlide], ordinal: int, budget: ChunkBudget
) -> list[Chunk]:
    display = _slide_text(slides)
    first, last = slides[0].number, slides[-1].number
    boundaries = tuple(f"slide:{slide.number}" for slide in slides)
    coordinates = ChunkCoordinates(
        heading_path=tuple(
            item for item in (slides[0].deck_title, slides[0].section, slides[0].title) if item
        ),
        slide=first if first == last else None,
        slide_start=first,
        slide_end=last,
        internal_boundaries=boundaries,
    )
    if len(display) <= budget.max_display_chars:
        built = build_chunk(
            run_id=run_id,
            file_id=file_id,
            ordinal=ordinal,
            pipeline_type="presentation",
            chunk_type="presentation",
            display_text=display,
            vector_text=_retrieval_text(slides),
            fulltext=_retrieval_text(slides),
            coordinates=coordinates,
            budget=budget,
        )
        return built if isinstance(built, list) else [built]
    context = "\n".join(f"Slide {slide.number}: {slide.title}" for slide in slides)
    body = "\n\n".join(part for part in (_slide_text([slide]) for slide in slides) if part)
    # The slide breadcrumb is repeated for every oversized piece, while the
    # body is bounded against both retrieval fields so build_chunk can still
    # enforce the shared contract.
    separator = "\n\n"
    piece_limit = min(
        budget.max_display_chars - len(context) - len(separator),
        budget.max_vector_chars - len(context) - len(separator),
        budget.max_fulltext_chars - len(context) - len(separator),
    )
    piece_limit = max(1, piece_limit)
    body = body.strip()
    body = body if body else context
    words = body.split()
    pieces: list[str] = []
    current = ""
    for word in words:
        if current and len(current) + len(word) + 1 > piece_limit:
            pieces.append(current)
            current = word
        else:
            current = word if not current else f"{current} {word}"
    if current:
        pieces.append(current)
    chunks: list[Chunk] = []
    for index, piece in enumerate(pieces):
        repeated_context = f"{context}{separator}{piece}"
        built = build_chunk(
            run_id=run_id,
            file_id=file_id,
            ordinal=ordinal + index,
            pipeline_type="presentation",
            chunk_type="presentation",
            display_text=repeated_context,
            vector_text=repeated_context,
            fulltext=repeated_context,
            coordinates=coordinates,
            budget=budget,
        )
        chunks.extend(built if isinstance(built, list) else [built])
    return chunks


def _slide_text(slides: list[PresentationSlide]) -> str:
    parts: list[str] = []
    for slide in slides:
        heading = f"Slide {slide.number}: {slide.title}".rstrip(": ")
        body = [block.text for block in slide.blocks]
        if slide.notes:
            body.append(f"Speaker notes:\n{slide.notes}")
        parts.append("\n".join([heading, *body]) if body else heading)
    return "\n\n".join(parts)


def _retrieval_text(slides: list[PresentationSlide]) -> str:
    context = "\n".join(f"Slide {slide.number}: {slide.title}" for slide in slides)
    return f"{context}\n\n{_slide_text(slides)}" if context else _slide_text(slides)


def _section_for_slide(title: str, number: int) -> str:
    lowered = title.lower()
    if number == 1 or lowered.startswith(("agenda", "contents", "introduction")):
        return title or "Introduction"
    return "default"


def _table_id(shape: Any) -> str:
    values = []
    for row in shape.table.rows:
        values.append("|".join(cell.text.strip() for cell in row.cells))
    return "toast_tbl_" + hashlib.sha256("\n".join(values).encode()).hexdigest()[:20]
