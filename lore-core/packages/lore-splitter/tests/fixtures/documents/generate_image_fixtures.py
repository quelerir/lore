from __future__ import annotations

import hashlib
from dataclasses import dataclass
from pathlib import Path

from docx import Document
from docx.shared import Inches
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches as PptxInches

try:
    import pymupdf as fitz
except ImportError:  # pragma: no cover - older PyMuPDF import path.
    import fitz


@dataclass(frozen=True)
class ImageContractFixtures:
    docx: Path
    pptx: Path
    pdf: Path
    scanned_pdf: Path
    content_image: Path
    decorative_logo: Path


def create_image_contract_fixtures(root: Path) -> ImageContractFixtures:
    root.mkdir(parents=True, exist_ok=True)
    content_image = root / "content-chart.png"
    decorative_logo = root / "decorative-logo.png"
    full_page_scan = root / "full-page-scan.png"
    docx = root / "image-policy.docx"
    pptx = root / "image-slides.pptx"
    pdf = root / "image-manual.pdf"
    scanned_pdf = root / "scanned-manual.pdf"

    _make_png(content_image, (320, 180), "Quarterly chart", (20, 90, 180))
    _make_png(decorative_logo, (42, 42), "L", (160, 40, 40))
    _make_png(full_page_scan, (1200, 1600), "Full page raster scan", (70, 70, 70))
    _write_docx(docx, content_image, decorative_logo)
    _write_pptx(pptx, content_image, decorative_logo)
    _write_pdf(pdf, content_image, decorative_logo)
    _write_scanned_pdf(scanned_pdf, full_page_scan)

    return ImageContractFixtures(
        docx=docx,
        pptx=pptx,
        pdf=pdf,
        scanned_pdf=scanned_pdf,
        content_image=content_image,
        decorative_logo=decorative_logo,
    )


def file_sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _make_png(path: Path, size: tuple[int, int], label: str, color: tuple[int, int, int]) -> None:
    image = Image.new("RGB", size, "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((4, 4, size[0] - 5, size[1] - 5), outline=color, width=4)
    draw.line((12, size[1] - 16, size[0] // 2, size[1] // 3), fill=color, width=5)
    draw.line((size[0] // 2, size[1] // 3, size[0] - 12, size[1] - 30), fill=color, width=5)
    draw.text((16, 16), label, fill=color)
    image.save(path)


def _write_docx(path: Path, content_image: Path, decorative_logo: Path) -> None:
    document = Document()
    document.add_heading("Image Policy", level=1)
    document.add_paragraph("The chart below is a true content image.")
    document.add_picture(str(content_image), width=Inches(3.2))
    document.add_paragraph("Repeated decorative marks should be skipped.")
    for _ in range(3):
        document.add_picture(str(decorative_logo), width=Inches(0.28))
    document.save(path)


def _write_pptx(path: Path, content_image: Path, decorative_logo: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Image Slide"
    slide.shapes.add_picture(
        str(content_image),
        PptxInches(0.8),
        PptxInches(1.4),
        width=PptxInches(3.2),
    )
    for index in range(3):
        slide.shapes.add_picture(
            str(decorative_logo),
            PptxInches(0.8 + index * 0.45),
            PptxInches(4.1),
            width=PptxInches(0.28),
        )
    presentation.save(path)


def _write_pdf(path: Path, content_image: Path, decorative_logo: Path) -> None:
    doc = fitz.open()
    for page_number in range(2):
        page = doc.new_page(width=612, height=792)
        page.insert_text((72, 72), f"PDF image fixture page {page_number + 1}")
        page.insert_image(fitz.Rect(72, 137, 392, 317), filename=str(content_image))
        if page_number == 0:
            for index in range(3):
                x0 = 72 + index * 54
                page.insert_image(fitz.Rect(x0, 354, x0 + 28, 382), filename=str(decorative_logo))
    doc.save(path)
    doc.close()


def _write_scanned_pdf(path: Path, full_page_scan: Path) -> None:
    doc = fitz.open()
    page = doc.new_page(width=612, height=792)
    page.insert_image(fitz.Rect(0, 0, 612, 792), filename=str(full_page_scan))
    page.insert_text((72, 72), "Scanned page with text marker")
    doc.save(path)
    doc.close()
